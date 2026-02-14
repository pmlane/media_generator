#!/usr/bin/env python3
"""
Generate an editable PowerPoint file from a background image + text layout.

Usage:
    python generate-pptx.py <background_image> <layout_json> <output_pptx>

The layout JSON contains width/height (pixels) and an array of text elements
with position, font, size, color, and anchor information. Positions are
converted from pixels to inches at 300 DPI.

Requires: python-pptx (pip install python-pptx)
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


DPI = 300


def px_to_inches(px: float) -> float:
    return px / DPI


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def get_alignment(anchor: str):
    if anchor == "middle":
        return PP_ALIGN.CENTER
    elif anchor == "end":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def main():
    if len(sys.argv) != 4:
        print(f"Usage: {sys.argv[0]} <background_image> <layout_json> <output_pptx>", file=sys.stderr)
        sys.exit(1)

    bg_path = Path(sys.argv[1])
    layout_path = Path(sys.argv[2])
    output_path = Path(sys.argv[3])

    if not bg_path.exists():
        print(f"Error: Background image not found: {bg_path}", file=sys.stderr)
        sys.exit(1)

    with open(layout_path) as f:
        layout = json.load(f)

    width_in = px_to_inches(layout["width"])
    height_in = px_to_inches(layout["height"])

    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background image filling the slide
    slide.shapes.add_picture(
        str(bg_path),
        Inches(0),
        Inches(0),
        Inches(width_in),
        Inches(height_in),
    )

    # Add text boxes for each element
    for el in layout["elements"]:
        x_px = el["x"]
        y_px = el["y"]
        font_size = el["fontSize"]
        max_width = el.get("maxWidth", layout["width"] * 0.8)
        anchor = el.get("anchor", "start")

        # Convert to inches
        box_width = px_to_inches(max_width)
        box_height = px_to_inches(font_size * 2)
        y_in = px_to_inches(y_px - font_size)  # SVG y is baseline; shift up

        # Adjust x based on anchor
        if anchor == "middle":
            x_in = px_to_inches(x_px) - box_width / 2
        elif anchor == "end":
            x_in = px_to_inches(x_px) - box_width
        else:
            x_in = px_to_inches(x_px)

        # Clamp to slide bounds
        x_in = max(0, x_in)
        y_in = max(0, y_in)

        txBox = slide.shapes.add_textbox(
            Inches(x_in),
            Inches(y_in),
            Inches(box_width),
            Inches(box_height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = el["text"]
        p.alignment = get_alignment(anchor)

        run = p.runs[0]
        run.font.size = Pt(font_size * 72 / DPI)  # Convert px at 300 DPI to points
        run.font.bold = el.get("fontWeight") == "bold"
        run.font.color.rgb = hex_to_rgb(el.get("color", "#000000"))

        # Set font family
        font_family = el.get("fontFamily", "Arial")
        run.font.name = font_family

    prs.save(str(output_path))
    print(str(output_path))


if __name__ == "__main__":
    main()
